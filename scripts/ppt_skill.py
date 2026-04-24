from __future__ import annotations

import math
import zipfile
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt as PptPt


SLIDE_W = 13.333
SLIDE_H = 7.5


PALETTE = {
    "bg": RGBColor(255, 255, 255),
    "ink": RGBColor(31, 41, 55),
    "muted": RGBColor(107, 114, 128),
    "rule": RGBColor(209, 213, 219),
    "accent": RGBColor(37, 99, 235),
    "accent_soft": RGBColor(219, 234, 254),
    "amber": RGBColor(180, 83, 9),
    "amber_soft": RGBColor(254, 243, 199),
    "green": RGBColor(21, 128, 61),
    "green_soft": RGBColor(220, 252, 231),
    "code_bg": RGBColor(17, 24, 39),
    "code_fg": RGBColor(229, 231, 235),
    "output_bg": RGBColor(240, 253, 244),
}


def add_text(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    size: int,
    *,
    bold: bool = False,
    color: RGBColor | None = None,
    align=PP_ALIGN.LEFT,
    font_name="Microsoft YaHei",
    fill: RGBColor | None = None,
    line: RGBColor | None = None,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    if line is not None:
        box.line.color.rgb = line
    else:
        box.line.fill.background()
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    for run in paragraph.runs:
        run.font.name = font_name
        run.font.size = PptPt(size)
        run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color
    return box


def add_rect(slide, left: float, top: float, width: float, height: float, fill: RGBColor, line: RGBColor | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def add_round_rect(slide, left: float, top: float, width: float, height: float, fill: RGBColor, line: RGBColor | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def set_text_frame(shape, text: str, size: int, *, color: RGBColor, font_name="Microsoft YaHei", bold: bool = False) -> None:
    frame = shape.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.14)
    frame.margin_right = Inches(0.14)
    frame.margin_top = Inches(0.08)
    frame.margin_bottom = Inches(0.08)
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    for run in paragraph.runs:
        run.font.name = font_name
        run.font.size = PptPt(size)
        run.font.color.rgb = color
        run.font.bold = bold


def add_bullets(slide, items: list[str], left: float, top: float, width: float, height: float, size: int = 15):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.12)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)
    frame.clear()
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.space_after = PptPt(5)
        paragraph.font.name = "Microsoft YaHei"
        paragraph.font.size = PptPt(size)
        paragraph.font.color.rgb = PALETTE["ink"]
    return box


def add_code_block(slide, code: str, left: float, top: float, width: float, height: float):
    shape = add_round_rect(slide, left, top, width, height, PALETTE["code_bg"], RGBColor(55, 65, 81))
    frame = shape.text_frame
    frame.word_wrap = False
    frame.margin_left = Inches(0.18)
    frame.margin_right = Inches(0.12)
    frame.margin_top = Inches(0.12)
    frame.margin_bottom = Inches(0.10)
    frame.clear()
    keywords = {"for", "in", "def", "return", "import"}
    for line_index, line in enumerate(code.rstrip().splitlines()):
        paragraph = frame.paragraphs[0] if line_index == 0 else frame.add_paragraph()
        paragraph.font.name = "Consolas"
        paragraph.font.size = PptPt(10.5)
        parts = line.replace("(", " (").replace(")", ") ").replace(":", ": ").split(" ")
        for part in parts:
            run = paragraph.add_run()
            run.text = part + (" " if part else "")
            run.font.name = "Consolas"
            run.font.size = PptPt(10.5)
            run.font.color.rgb = RGBColor(147, 197, 253) if part in keywords else PALETTE["code_fg"]
    return shape


def formula_lookup(slide_data: dict) -> dict[str, dict]:
    return {item["formula_id"]: item for item in slide_data["formula_spec"]}


def block_lookup(slide_data: dict) -> dict[str, dict]:
    return {item["block_id"]: item for item in slide_data["content_blocks"]}


def generate_loss_curve(asset_path: Path) -> None:
    x = np.linspace(-2, 8, 400)
    y = (x - 3) ** 2 + 1
    w0 = 0
    tangent_x = np.linspace(-0.8, 2.3, 100)
    tangent_y = (w0 - 3) ** 2 + 1 + 2 * (w0 - 3) * (tangent_x - w0)
    plt.figure(figsize=(6.2, 3.7))
    plt.plot(x, y, color="#2563eb", linewidth=2.5, label="L(w)")
    plt.plot(tangent_x, tangent_y, color="#f59e0b", linestyle="--", linewidth=2, label="local slope")
    plt.scatter([w0, 3], [(w0 - 3) ** 2 + 1, 1], color=["#dc2626", "#16a34a"], s=55)
    plt.annotate("move right", xy=(0.6, 8), xytext=(1.8, 14), arrowprops={"arrowstyle": "->", "color": "#dc2626"})
    plt.title("Toy loss curve")
    plt.xlabel("parameter w")
    plt.ylabel("loss L(w)")
    plt.grid(alpha=0.22)
    plt.legend()
    plt.tight_layout()
    plt.savefig(asset_path, dpi=180)
    plt.close()


def loss(w):
    return (w - 3) ** 2 + 1


def grad(w):
    return 2 * (w - 3)


def run_descent(eta: float, steps: int = 25, w0: float = 0.0) -> tuple[np.ndarray, np.ndarray]:
    ws = [w0]
    losses = [loss(w0)]
    w = w0
    for _ in range(steps):
        w = w - eta * grad(w)
        ws.append(w)
        losses.append(loss(w))
    return np.array(ws), np.array(losses)


def generate_learning_rate_trajectories(asset_path: Path) -> None:
    learning_rates = [0.05, 0.2, 1.05]
    results = {eta: run_descent(eta) for eta in learning_rates}
    fig, axes = plt.subplots(1, 2, figsize=(7.4, 3.5))
    x = np.linspace(-2, 8, 300)
    axes[0].plot(x, loss(x), color="#111827", linewidth=2, label="L(w)")
    for eta, (ws, losses) in results.items():
        axes[0].plot(ws, loss(ws), marker="o", markersize=3, linewidth=1.5, label=f"eta={eta}")
        axes[1].plot(losses, marker="o", markersize=3, linewidth=1.5, label=f"eta={eta}")
    axes[0].set_title("parameter trajectory")
    axes[0].set_xlabel("w")
    axes[0].set_ylabel("L(w)")
    axes[1].set_title("loss by step")
    axes[1].set_xlabel("step")
    axes[1].set_ylabel("loss")
    for axis in axes:
        axis.grid(alpha=0.22)
        axis.legend(fontsize=8)
    plt.tight_layout()
    plt.savefig(asset_path, dpi=180)
    plt.close()


def generate_box_diagram(asset_path: Path, labels: list[str], title: str) -> None:
    fig, ax = plt.subplots(figsize=(6.8, 3.6))
    ax.set_axis_off()
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 6)
    coords = [(1, 3.2), (3.4, 3.2), (5.8, 3.2), (8.2, 3.2)]
    for index, (x, y) in enumerate(coords):
        rect = plt.Rectangle((x - 0.85, y - 0.55), 1.7, 1.1, facecolor="#dbeafe", edgecolor="#2563eb", linewidth=1.8)
        ax.add_patch(rect)
        ax.text(x, y, labels[index], ha="center", va="center", fontsize=11, color="#1f2937")
        if index < len(coords) - 1:
            ax.annotate("", xy=(coords[index + 1][0] - 0.95, y), xytext=(x + 0.95, y), arrowprops={"arrowstyle": "->", "lw": 2, "color": "#4b5563"})
    ax.text(5, 5.25, title, ha="center", va="center", fontsize=14, weight="bold", color="#111827")
    plt.tight_layout()
    plt.savefig(asset_path, dpi=180, transparent=False)
    plt.close()


def generate_assets(output_path: Path) -> dict[str, Path]:
    asset_dir = output_path.parent / "assets"
    asset_dir.mkdir(parents=True, exist_ok=True)
    assets = {
        "loss_curve": asset_dir / "loss_curve.png",
        "learning_rate_trajectories": asset_dir / "learning_rate_trajectories.png",
        "training_loop": asset_dir / "training_loop.png",
        "parameter_update": asset_dir / "parameter_update.png",
    }
    generate_loss_curve(assets["loss_curve"])
    generate_learning_rate_trajectories(assets["learning_rate_trajectories"])
    generate_box_diagram(assets["training_loop"], ["predict", "loss", "gradient", "update"], "training loop")
    generate_box_diagram(assets["parameter_update"], ["loss", "gradient", "eta", "theta"], "error becomes an update")
    return assets


def add_slide_chrome(slide, slide_data: dict, palette: dict[str, RGBColor], *, footer: str) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = palette["bg"]
    add_text(slide, slide_data["title"], 0.45, 0.24, 11.8, 0.55, 23, bold=True, color=palette["ink"])
    add_text(slide, slide_data["teaching_purpose"], 0.48, 0.82, 11.8, 0.34, 10, color=palette["muted"])
    add_rect(slide, 0.48, 1.14, 12.25, 0.02, palette["rule"])
    add_text(slide, footer, 10.2, 7.08, 2.5, 0.25, 8, color=palette["muted"], align=PP_ALIGN.RIGHT)


def render_block(slide, slide_data: dict, block: dict, x: float, y: float, w: float, h: float, assets: dict[str, Path]):
    kind = block["kind"]
    formulas = formula_lookup(slide_data)
    if kind == "paragraph":
        return add_text(slide, block["text"], x, y, w, h, 14, color=PALETTE["ink"])
    if kind == "bullet_list":
        return add_bullets(slide, block["items"], x, y, w, h)
    if kind == "callout":
        shape = add_round_rect(slide, x, y, w, h, PALETTE["amber_soft"], RGBColor(245, 158, 11))
        set_text_frame(shape, block["text"], 15, color=PALETTE["amber"], bold=True)
        return shape
    if kind == "formula":
        spec = formulas[block["formula_ref"]]
        formula_text = f"\\[{spec['latex']}\\]\n{spec['spoken_explanation']}"
        shape = add_round_rect(slide, x, y, w, h, PALETTE["accent_soft"], PALETTE["accent"])
        set_text_frame(shape, formula_text, 15, color=PALETTE["ink"], font_name="Cambria Math")
        return shape
    if kind == "code_block":
        return add_code_block(slide, block["code"], x, y, w, h)
    if kind == "output_block":
        shape = add_round_rect(slide, x, y, w, h, PALETTE["output_bg"], PALETTE["green"])
        set_text_frame(shape, block["text"], 13, color=PALETTE["green"], bold=True)
        return shape
    if kind in {"diagram", "plot"}:
        image = slide.shapes.add_picture(str(assets[block["asset_id"]]), Inches(x), Inches(y), width=Inches(w), height=Inches(h - 0.32))
        caption = add_text(slide, block.get("caption", ""), x, y + h - 0.28, w, 0.25, 9, color=PALETTE["muted"], align=PP_ALIGN.CENTER)
        return image
    raise ValueError(f"Unsupported content block kind: {kind}")


def arrange_blocks(slide_data: dict) -> list[tuple[dict, float, float, float, float]]:
    blocks = slide_data["content_blocks"]
    layout = slide_data["layout"]
    left_x, right_x = 0.62, 6.92
    left_w, right_w = 5.85, 5.75
    if layout in {"title_overview", "two_column", "diagram_explanation"}:
        left_blocks = [b for b in blocks if b["kind"] in {"paragraph", "bullet_list", "callout", "formula"}]
        right_blocks = [b for b in blocks if b["kind"] in {"diagram", "plot", "code_block", "output_block"}]
        placements = []
        y = 1.38
        for block in left_blocks:
            h = 1.05 if block["kind"] == "callout" else 1.32 if block["kind"] == "paragraph" else 1.55
            placements.append((block, left_x, y, left_w, h))
            y += h + 0.18
        visual_h = 4.75 if len(right_blocks) == 1 else 2.35
        y = 1.45
        for block in right_blocks:
            placements.append((block, right_x, y, right_w, visual_h))
            y += visual_h + 0.32
        return placements
    if layout == "formula_lab":
        placements = []
        y = 1.34
        for block in blocks:
            if block["kind"] == "paragraph":
                placements.append((block, 0.65, y, 12.05, 0.95))
                y += 1.1
        formula_blocks = [b for b in blocks if b["kind"] == "formula"]
        for index, block in enumerate(formula_blocks):
            placements.append((block, 0.8 + index * 5.95, y, 5.6, 1.45))
        for block in blocks:
            if block["kind"] == "callout":
                placements.append((block, 2.0, 5.1, 9.35, 0.95))
        return placements
    if layout == "code_result":
        placements = []
        for block in blocks:
            if block["kind"] == "paragraph":
                placements.append((block, 0.65, 1.32, 5.9, 0.9))
            elif block["kind"] == "code_block":
                placements.append((block, 0.65, 1.55, 5.9, 2.35))
            elif block["kind"] == "output_block":
                placements.append((block, 0.65, 4.1, 5.9, 1.0))
            elif block["kind"] == "formula":
                placements.append((block, 0.65, 2.5, 5.9, 1.25))
            elif block["kind"] in {"plot", "diagram"}:
                placements.append((block, 6.85, 1.48, 5.85, 4.65))
        return placements
    raise ValueError(f"Unsupported layout: {layout}")


def visible_blocks_for_reveal(slide_data: dict, reveal_index: int | None) -> set[str]:
    if reveal_index is None:
        return {block["block_id"] for block in slide_data["content_blocks"]}
    return {step["target_ref"] for step in slide_data["reveal_steps"][: reveal_index + 1]}


def render_logical_slide(slide, slide_data: dict, assets: dict[str, Path], *, footer: str, reveal_index: int | None = None) -> dict[str, int]:
    add_slide_chrome(slide, slide_data, PALETTE, footer=footer)
    visible_ids = visible_blocks_for_reveal(slide_data, reveal_index)
    shape_ids: dict[str, int] = {}
    for block, x, y, w, h in arrange_blocks(slide_data):
        if block["block_id"] not in visible_ids:
            continue
        shape = render_block(slide, slide_data, block, x, y, w, h, assets)
        shape_ids[block["block_id"]] = shape.shape_id
    add_notes(slide, slide_data["notes"])
    return shape_ids


def add_notes(slide, notes: str) -> None:
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = notes


def animation_xml(shape_ids: list[int]) -> str:
    child_nodes = []
    next_id = 3
    for shape_id in shape_ids:
        child_nodes.append(
            f"""
            <p:par>
              <p:cTn id="{next_id}" fill="hold">
                <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
                <p:childTnLst>
                  <p:par>
                    <p:cTn id="{next_id + 1}" presetID="1" presetClass="entr" presetSubtype="0" fill="hold" grpId="0" nodeType="clickEffect">
                      <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                      <p:childTnLst>
                        <p:set>
                          <p:cBhvr>
                            <p:cTn id="{next_id + 2}" dur="1" fill="hold"/>
                            <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                            <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                          </p:cBhvr>
                          <p:to><p:strVal val="visible"/></p:to>
                        </p:set>
                      </p:childTnLst>
                    </p:cTn>
                  </p:par>
                </p:childTnLst>
              </p:cTn>
            </p:par>
            """
        )
        next_id += 3
    return f"""
    <p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
      <p:tnLst>
        <p:par>
          <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
            <p:childTnLst>
              <p:seq concurrent="1" nextAc="seek">
                <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
                  <p:childTnLst>
                    {''.join(child_nodes)}
                  </p:childTnLst>
                </p:cTn>
                <p:prevCondLst><p:cond evt="onPrev" delay="0"/></p:prevCondLst>
                <p:nextCondLst><p:cond evt="onNext" delay="0"/></p:nextCondLst>
              </p:seq>
            </p:childTnLst>
          </p:cTn>
        </p:par>
      </p:tnLst>
    </p:timing>
    """


def inject_reveal_timing(slide, ordered_shape_ids: list[int]) -> None:
    if not ordered_shape_ids:
        return
    timing = parse_xml(animation_xml(ordered_shape_ids))
    slide._element.append(timing)


def presentation_has_timing(path: Path) -> bool:
    with zipfile.ZipFile(path) as archive:
        slide_names = [name for name in archive.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")]
        for name in slide_names:
            xml = archive.read(name).decode("utf-8", errors="ignore")
            if "<p:timing" in xml:
                return True
    return False


def build_presentation(slide_script: dict, assets: dict[str, Path], *, physical_reveal: bool) -> tuple[Presentation, dict]:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slides = slide_script["slide_script"]["slides"]
    native_steps = 0
    physical_pages = 0
    for index, slide_data in enumerate(slides, start=1):
        if physical_reveal:
            for reveal_index in range(len(slide_data["reveal_steps"])):
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                render_logical_slide(
                    slide,
                    slide_data,
                    assets,
                    footer=f"{slide_data['slide_id']} · reveal {reveal_index + 1}/{len(slide_data['reveal_steps'])}",
                    reveal_index=reveal_index,
                )
                physical_pages += 1
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape_ids = render_logical_slide(slide, slide_data, assets, footer=f"{slide_data['slide_id']} · {index}/{len(slides)}")
        ordered_ids = [shape_ids[step["target_ref"]] for step in slide_data["reveal_steps"] if step["target_ref"] in shape_ids]
        inject_reveal_timing(slide, ordered_ids)
        native_steps += len(ordered_ids)
        physical_pages += 1
    return prs, {"native_reveal_steps": native_steps, "physical_pages": physical_pages}


def render_pptx(slide_script: dict, visual_config: dict, output_path: Path) -> None:
    assets = generate_assets(output_path)
    render_report = {
        "mode": "native_animation",
        "fallback_used": False,
        "asset_count": len(assets),
        "generated_assets": {key: str(path) for key, path in assets.items()},
    }
    try:
        prs, stats = build_presentation(slide_script, assets, physical_reveal=False)
        prs.save(output_path)
        if not presentation_has_timing(output_path):
            raise RuntimeError("Native animation XML was not found after saving.")
        render_report.update(stats)
    except PermissionError:
        raise
    except Exception as exc:
        prs, stats = build_presentation(slide_script, assets, physical_reveal=True)
        prs.save(output_path)
        render_report.update(stats)
        render_report["mode"] = "physical_reveal_fallback"
        render_report["fallback_used"] = True
        render_report["fallback_reason"] = str(exc)
    slide_script["slide_script"]["render_report"] = render_report


def write_ppt_skill_report(slide_script: dict, visual_config: dict, output_path: Path) -> None:
    root = slide_script["slide_script"]
    visual = visual_config["visual_ppt"]
    verdict = root["visual_verdict"]
    render_report = root.get("render_report", {})
    lines = [
        f"# {root['title']} PPT Skill 报告",
        "",
        "## 职责边界",
        "",
        f"- {visual['purpose']}",
    ]
    lines.extend([f"- {item}" for item in visual["boundary"]])
    lines.extend([
        "",
        "## 讲义型渲染",
        "",
        f"- 逻辑页：{len(root['slides'])}",
        f"- 内容块：{sum(len(slide['content_blocks']) for slide in root['slides'])}",
        f"- Reveal 步骤：{sum(len(slide['reveal_steps']) for slide in root['slides'])}",
        f"- Reveal 实现：{render_report.get('mode', 'unknown')}",
        f"- 物理页：{render_report.get('physical_pages', 'unknown')}",
        f"- 生成资产：{render_report.get('asset_count', 0)}",
        f"- 兜底启用：{render_report.get('fallback_used', False)}",
        "",
        "## 视觉质量",
        "",
        f"- 结论：{verdict['verdict']}",
        f"- 总分：{verdict['total_score']} / 100",
        f"- 通过阈值：{verdict['pass_threshold']}",
        "",
        "## 风险",
        "",
    ])
    lines.extend([f"- {item}" for item in verdict["risks"]])
    if render_report.get("fallback_reason"):
        lines.append(f"- Reveal 兜底原因：{render_report['fallback_reason']}")
    lines.extend([
        "",
        "## 修复动作",
        "",
    ])
    lines.extend([f"- {item}" for item in verdict["repair_actions"]])
    output_path.write_text("\n".join(lines) + "\n", encoding="utf-8")
